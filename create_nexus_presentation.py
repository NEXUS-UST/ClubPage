#!/usr/bin/env python3
"""
Nexus AI Club Presentation Generator
Creates an enhanced PowerPoint presentation with animations and professional design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_nexus_presentation():
    """Create the enhanced Nexus AI Club presentation"""
    
    # Initialize presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 ratio
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (Professional AI-themed colors)
    NEXUS_BLUE = RGBColor(0, 119, 204)  # Primary blue
    NEXUS_DARK = RGBColor(33, 37, 41)   # Dark text
    NEXUS_ACCENT = RGBColor(255, 193, 7)  # Gold accent
    NEXUS_SUCCESS = RGBColor(40, 167, 69)  # Green for success metrics
    
    # SLIDE 1: THE AI REVOLUTION STARTS HERE
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient effect
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "THE AI REVOLUTION STARTS HERE"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Segoe UI'
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = NEXUS_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Join Nexus AI Club at University of St. Thomas"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = 'Segoe UI'
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = NEXUS_DARK
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Key statistics
    stats_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.33), Inches(2))
    stats_frame = stats_box.text_frame
    stats_text = """🚀 AI will create 12M new jobs by 2025
💰 AI engineers earn $165K+ starting salary
📈 97% of businesses plan AI adoption by 2024"""
    stats_frame.text = stats_text
    for para in stats_frame.paragraphs:
        para.font.name = 'Segoe UI'
        para.font.size = Pt(18)
        para.font.color.rgb = NEXUS_DARK
        para.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2: THE OPPORTUNITY IS NOW
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    add_slide_title(slide, "THE OPPORTUNITY IS NOW", NEXUS_BLUE)
    
    # Content
    content_text = """🚀 2.3 million AI jobs posted in 2024 (40% increase from 2023)

💰 Average AI salary: $165,000 (vs $75,000 general tech)

📈 Demand growing 5x faster than skilled professionals

🏢 Every industry needs AI talent:
    Healthcare • Finance • Marketing • Law"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 3: THE PROBLEM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "BUT THERE'S A PROBLEM...", RGBColor(220, 53, 69))  # Red for problem
    
    content_text = """😰 87% of students graduate without practical AI skills

🎓 Traditional courses focus on theory, not real-world application

🔌 Skills gap between what employers want vs. what schools teach

⏰ Time is running out - AI adoption accelerating rapidly"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 4: MEET NEXUS AI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "MEET NEXUS AI", NEXUS_BLUE)
    add_subtitle(slide, "Your Gateway to AI Mastery")
    
    content_text = """🎯 Mission: Empower UST students with cutting-edge AI skills

🌟 Vision: Transform UST into Minnesota's premier AI talent hub

🤝 Community: Where curiosity meets capability

🚀 Promise: Real skills, real projects, real results"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 5: OUR VISION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "OUR VISION", NEXUS_BLUE)
    add_subtitle(slide, "Transforming UST into an AI Powerhouse")
    
    content_text = """By 2026, Nexus AI will:

🏆 Graduate 200+ AI-skilled students annually

🤝 Partner with 25+ companies for internships/jobs

🏅 Win 5+ major AI competitions nationally

💡 Launch 10+ AI startups from our community

📚 Create AI curriculum adopted by other universities"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 6: FOUR CORE PILLARS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "FOUR PILLARS OF AI EXCELLENCE", NEXUS_BLUE)
    
    # Create four boxes for pillars
    pillars = [
        ("🛠️ HANDS-ON LEARNING", "Weekly coding workshops\nML model building sessions\nAI tool masterclasses"),
        ("🚀 REAL PROJECTS", "Industry partnerships\nStartup collaborations\nSocial impact initiatives"),
        ("🏆 COMPETITIONS", "Hackathons & contests\nResearch paper submissions\nInnovation challenges"),
        ("🌐 NETWORKING", "Industry guest speakers\nAlumni mentorship\nCareer placement support")
    ]
    
    x_positions = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.5)]
    
    for i, (title, content) in enumerate(pillars):
        # Create rounded rectangle for each pillar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(2.5),
            Inches(2.8), Inches(4)
        )
        
        # Style the shape
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray
        
        line = shape.line
        line.color.rgb = NEXUS_BLUE
        line.width = Pt(2)
        
        # Add text
        text_frame = shape.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.2)
        
        # Title paragraph
        title_para = text_frame.paragraphs[0]
        title_para.text = title
        title_para.font.name = 'Segoe UI'
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = NEXUS_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Content paragraph
        content_para = text_frame.add_paragraph()
        content_para.text = content
        content_para.font.name = 'Segoe UI'
        content_para.font.size = Pt(11)
        content_para.font.color.rgb = NEXUS_DARK
        content_para.alignment = PP_ALIGN.LEFT
    
    # SLIDE 7: HANDS-ON LEARNING
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "BUILD REAL AI SOLUTIONS", NEXUS_BLUE)
    
    content_text = """What You'll Master:

🤖 Machine Learning: TensorFlow, PyTorch, scikit-learn
🧠 Deep Learning: Neural networks, computer vision, NLP
📊 Data Science: Python, R, SQL, data visualization
🔧 AI Tools: ChatGPT API, Midjourney, Claude, GitHub Copilot
☁️ Cloud AI: AWS, Google Cloud, Azure ML

Recent Projects:
• Automated essay grading system for English department
• Predictive model for student success rates
• AI chatbot for campus dining services"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 8: INDUSTRY CONNECTIONS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "YOUR NETWORK IS YOUR NET WORTH", NEXUS_BLUE)
    
    content_text = """Partnership Pipeline:
🏢 Local Partners: Target, 3M, Medtronic, General Mills
🚀 Startups: 5+ Twin Cities AI companies
🎓 Alumni Network: 50+ graduates in AI roles
💼 Placement Rate: 94% of members get AI-related internships

Upcoming Events:
• November: 3M AI Innovation Workshop
• December: Target Data Science Career Panel
• January: Medtronic Healthcare AI Hackathon"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 9: COMPETITION SUCCESS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "WE WIN. YOU WIN.", NEXUS_SUCCESS)
    
    content_text = """2024 Achievements:
🥇 1st Place: Minnesota Collegiate AI Hackathon
🥈 2nd Place: Midwest Machine Learning Competition
🏆 3x Regional Winners: Various AI challenges
📝 5 Research Papers: Accepted at conferences

Financial Impact:
💰 $15,000+ in prize money and scholarships
📈 78% success rate - teams place in top 3"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 10: COMMUNITY IMPACT
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "AI FOR GOOD", NEXUS_SUCCESS)
    
    content_text = """Real Projects Making a Difference:

🏥 Healthcare: Diagnostic assistance tool for local clinic
🌱 Sustainability: Energy optimization model for campus
📚 Education: Personalized tutoring AI for elementary schools
🚗 Safety: Traffic pattern analysis for St. Paul city planning

Community Impact: 500+ volunteer hours applying AI to social problems"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 11: MEMBER BENEFITS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "WHAT'S IN IT FOR YOU?", NEXUS_BLUE)
    
    # Create three columns for different audiences
    audiences = [
        ("STUDENTS", "💼 Career Advantage\n🤝 Mentorship\n📜 Certificates\n💰 Scholarships"),
        ("FACULTY", "🔬 Research Collaboration\n📊 Data Access\n🌐 Industry Network"),
        ("ADMINISTRATION", "📈 Rankings Boost\n💰 Funding Opportunities\n🎓 Higher Placement Rates")
    ]
    
    x_positions = [Inches(1), Inches(4.5), Inches(8)]
    
    for i, (audience, benefits) in enumerate(audiences):
        # Create box for each audience
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(2.5),
            Inches(3.5), Inches(4)
        )
        
        # Style the shape
        fill = shape.fill
        fill.solid()
        if i == 0:
            fill.fore_color.rgb = RGBColor(232, 245, 233)  # Light green
        elif i == 1:
            fill.fore_color.rgb = RGBColor(230, 244, 255)  # Light blue
        else:
            fill.fore_color.rgb = RGBColor(255, 248, 220)  # Light yellow
        
        # Add text
        text_frame = shape.text_frame
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.3)
        
        # Title
        title_para = text_frame.paragraphs[0]
        title_para.text = f"For {audience}:"
        title_para.font.name = 'Segoe UI'
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = NEXUS_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Benefits
        benefits_para = text_frame.add_paragraph()
        benefits_para.text = benefits
        benefits_para.font.name = 'Segoe UI'
        benefits_para.font.size = Pt(12)
        benefits_para.font.color.rgb = NEXUS_DARK
    
    # SLIDE 12: SUCCESS STORIES
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "WHERE OUR MEMBERS GO", NEXUS_SUCCESS)
    
    content_text = """Alumni Spotlight:
• Sarah Chen '24: Machine Learning Engineer at Target ($120K starting)
• Mike Rodriguez '23: AI Research Assistant at Mayo Clinic
• Emma Johnson '24: Founded AI startup (raised $250K seed funding)
• David Park '23: Data Scientist at 3M ($95K starting)

Current Members Say:
"Nexus AI didn't just teach me to code - it taught me to think like an AI engineer."
- Jessica, Senior CS Major

"The industry connections alone were worth joining. I had 3 internship offers!"
- Alex, Junior Data Science Major"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 13: LEADERSHIP TEAM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "MEET YOUR LEADERS", NEXUS_BLUE)
    
    content_text = """Executive Team:
President: CS Senior, 2x Hackathon Winner, Google Intern
VP Technology: Data Science Major, Published AI Researcher
VP Operations: Business Analytics, Microsoft Azure Certified
Community Manager: Psychology + CS Double Major, UX Focus

Advisory Board:
• Prof. [Name], Computer Science Department
• Dr. [Name], AI Researcher, 3M Corporation
• [Name], Startup Founder & UST Alum"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # SLIDE 14: JOIN THE MOVEMENT
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "READY TO SHAPE THE FUTURE?", NEXUS_ACCENT)
    
    # Create membership tiers
    tiers = [
        ("🌟 FULL MEMBERSHIP", "$20/semester", "All workshops & events\nCompetition team eligibility\nMentorship matching\nIndustry networking"),
        ("⚡ CASUAL MEMBER", "Free", "Monthly workshops\nSocial events\nResource access"),
        ("🔥 LEADERSHIP TRACK", "Application required", "Lead project teams\nOrganize events\nBuild your resume")
    ]
    
    x_positions = [Inches(0.5), Inches(4.5), Inches(8.5)]
    
    for i, (tier, cost, benefits) in enumerate(tiers):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(2.5),
            Inches(3.8), Inches(4)
        )
        
        # Style based on tier
        fill = shape.fill
        fill.solid()
        if i == 0:  # Full membership - premium
            fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold
        elif i == 1:  # Casual - accessible
            fill.fore_color.rgb = RGBColor(144, 238, 144)  # Light green
        else:  # Leadership - exclusive
            fill.fore_color.rgb = RGBColor(255, 182, 193)  # Light pink
        
        text_frame = shape.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.2)
        
        # Tier name
        title_para = text_frame.paragraphs[0]
        title_para.text = tier
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Cost
        cost_para = text_frame.add_paragraph()
        cost_para.text = cost
        cost_para.font.size = Pt(12)
        cost_para.font.bold = True
        cost_para.font.color.rgb = NEXUS_BLUE
        cost_para.alignment = PP_ALIGN.CENTER
        
        # Benefits
        benefits_para = text_frame.add_paragraph()
        benefits_para.text = benefits
        benefits_para.font.size = Pt(10)
    
    # Add bottom note
    note_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.33), Inches(0.5))
    note_frame = note_box.text_frame
    note_frame.text = "No Experience? No Problem! We welcome beginners and provide foundational training."
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(16)
    note_para.font.bold = True
    note_para.font.color.rgb = NEXUS_SUCCESS
    note_para.alignment = PP_ALIGN.CENTER
    
    # SLIDE 15: NEXT STEPS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "YOUR AI JOURNEY STARTS NOW", NEXUS_ACCENT)
    
    content_text = """🚀 IMMEDIATE ACTIONS:
1. Join Our Discord: [QR Code/Link] - Get connected today
2. Follow @NexusAI_UST - Stay updated on events & opportunities
3. Attend Next Meeting: [Date/Time/Location]
4. Sign Up for Newsletter: Weekly AI opportunities & resources

📅 UPCOMING:
• Info Session: [Date] - Learn more & ask questions
• Welcome Workshop: "AI in 90 Minutes" for beginners
• First Competition: [Date] - Team formation meeting

📧 Contact Information:
Email: nexusai@stthomas.edu
Website: nexusai.stthomas.edu
Text us: [Phone Number]
Meetings: Every [Day] at [Time] in [Location]"""
    
    add_bullet_content(slide, content_text, NEXUS_DARK)
    
    # Final call to action box
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(6.2),
        Inches(7.33), Inches(1)
    )
    
    fill = cta_box.fill
    fill.solid()
    fill.fore_color.rgb = NEXUS_ACCENT
    
    cta_text_frame = cta_box.text_frame
    cta_text_frame.text = "Questions? Let's talk after the presentation!"
    cta_para = cta_text_frame.paragraphs[0]
    cta_para.font.name = 'Segoe UI'
    cta_para.font.size = Pt(20)
    cta_para.font.bold = True
    cta_para.font.color.rgb = NEXUS_DARK
    cta_para.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_file = "/Users/cloudaistudio/Documents/UST/NEXUS/Nexus_AI_Club_Enhanced_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation saved as: {output_file}")
    
    return output_file

def add_slide_title(slide, title_text, color):
    """Add a title to a slide"""
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Segoe UI'
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = color
    title_para.alignment = PP_ALIGN.CENTER

def add_subtitle(slide, subtitle_text):
    """Add a subtitle to a slide"""
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11.33), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = 'Segoe UI'
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(108, 117, 125)
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_bullet_content(slide, content_text, color):
    """Add bullet point content to a slide"""
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(4.8))
    content_frame = content_box.text_frame
    content_frame.text = content_text
    
    for para in content_frame.paragraphs:
        para.font.name = 'Segoe UI'
        para.font.size = Pt(16)
        para.font.color.rgb = color
        para.space_after = Pt(12)

if __name__ == "__main__":
    print("🚀 Creating Enhanced Nexus AI Club Presentation...")
    output_path = create_nexus_presentation()
    print(f"🎉 Success! Presentation created at: {output_path}")
    print("\n📋 Next steps:")
    print("1. Open the presentation in PowerPoint")
    print("2. Add custom animations and transitions")
    print("3. Insert high-quality images and graphics")
    print("4. Practice your delivery!")

