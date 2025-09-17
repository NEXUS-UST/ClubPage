#!/usr/bin/env python3
"""
Beautiful Nexus AI Club Presentation Creator
Creates a visually stunning PowerPoint with modern design principles
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_FILL
import math

class BeautifulNexusPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)  # 16:9 ratio
        self.prs.slide_height = Inches(9)
        
        # Professional color palette
        self.NEXUS_BLUE = RGBColor(0, 119, 204)
        self.NEXUS_DARK = RGBColor(33, 37, 41)
        self.NEXUS_ACCENT = RGBColor(255, 193, 7)
        self.NEXUS_SUCCESS = RGBColor(40, 167, 69)
        self.NEXUS_DANGER = RGBColor(220, 53, 69)
        self.NEXUS_LIGHT = RGBColor(248, 249, 250)
        self.NEXUS_GRADIENT_START = RGBColor(102, 126, 234)
        self.NEXUS_GRADIENT_END = RGBColor(118, 75, 162)
    
    def add_gradient_background(self, slide, color1, color2):
        """Add a gradient background to a slide"""
        try:
            # Create a rectangle that covers the entire slide
            background_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                self.prs.slide_width, self.prs.slide_height
            )
            
            # Set gradient fill
            fill = background_shape.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = color1
            fill.gradient_stops[1].color.rgb = color2
            
            # Send to back
            background_shape.element.getparent().insert(0, background_shape.element)
            
        except Exception as e:
            # Fallback to solid color if gradient fails
            fill = background_shape.fill
            fill.solid()
            fill.fore_color.rgb = color1
    
    def create_title_slide(self):
        """Create a stunning title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add gradient background
        self.add_gradient_background(slide, self.NEXUS_GRADIENT_START, self.NEXUS_GRADIENT_END)
        
        # Main title with shadow effect
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "THE AI REVOLUTION\nSTARTS HERE"
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        
        title_para = title_frame.paragraphs[0]
        title_para.font.name = 'Segoe UI Black'
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle with accent color
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Join Nexus AI Club at University of St. Thomas"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = 'Segoe UI'
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = self.NEXUS_ACCENT
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Create stunning stat cards
        stats = [
            ("12M", "New AI Jobs by 2025"),
            ("$165K+", "AI Engineer Starting Salary"),
            ("97%", "Businesses Adopting AI")
        ]
        
        card_width = Inches(4)
        card_height = Inches(1.5)
        start_x = Inches(2)
        y_pos = Inches(6.5)
        
        for i, (number, label) in enumerate(stats):
            x_pos = start_x + (i * Inches(4.5))
            
            # Create card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y_pos, card_width, card_height
            )
            
            # Style the card
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            fill.transparency = 0.8  # Semi-transparent
            
            # Add drop shadow effect
            shadow = card.shadow
            shadow.inherit = False
            shadow.style = 1  # Outer shadow
            shadow.blur_radius = Pt(10)
            shadow.distance = Pt(5)
            shadow.angle = 45
            
            # Add number
            num_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.1), card_width, Inches(0.7))
            num_frame = num_box.text_frame
            num_frame.text = number
            num_para = num_frame.paragraphs[0]
            num_para.font.name = 'Segoe UI Black'
            num_para.font.size = Pt(36)
            num_para.font.bold = True
            num_para.font.color.rgb = self.NEXUS_BLUE
            num_para.alignment = PP_ALIGN.CENTER
            
            # Add label
            label_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.8), card_width, Inches(0.6))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.name = 'Segoe UI'
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = self.NEXUS_DARK
            label_para.alignment = PP_ALIGN.CENTER
    
    def create_opportunity_slide(self):
        """Create the opportunity slide with visual impact"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Gradient background
        self.add_gradient_background(slide, RGBColor(74, 144, 226), RGBColor(180, 58, 255))
        
        # Title
        self.add_stunning_title(slide, "THE OPPORTUNITY IS NOW", RGBColor(255, 255, 255))
        
        # Create opportunity cards
        opportunities = [
            ("🚀", "2.3M AI Jobs", "Posted in 2024\n(40% increase)"),
            ("💰", "$165K Average", "AI Salary vs $75K\nGeneral Tech"),
            ("📈", "5x Faster", "Demand Growth\nvs Skilled Professionals"),
            ("🏢", "Every Industry", "Healthcare • Finance\nMarketing • Law")
        ]
        
        self.create_card_grid(slide, opportunities, start_y=Inches(2.5))
    
    def create_problem_slide(self):
        """Create the problem slide with urgency"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Red-tinted background for urgency
        self.add_gradient_background(slide, RGBColor(255, 99, 132), RGBColor(255, 159, 64))
        
        self.add_stunning_title(slide, "BUT THERE'S A PROBLEM...", RGBColor(255, 255, 255))
        
        problems = [
            ("😰", "87% of Students", "Graduate without\nPractical AI Skills"),
            ("🎓", "Traditional Courses", "Focus on Theory\nNot Application"),
            ("🔌", "Skills Gap", "Employer Needs vs\nSchool Teaching"),
            ("⏰", "Time Running Out", "AI Adoption\nAccelerating Rapidly")
        ]
        
        self.create_card_grid(slide, problems, start_y=Inches(2.5))
    
    def create_solution_slide(self):
        """Create the Nexus AI solution slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Success green gradient
        self.add_gradient_background(slide, RGBColor(56, 178, 172), RGBColor(129, 230, 217))
        
        self.add_stunning_title(slide, "MEET NEXUS AI", RGBColor(255, 255, 255))
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Your Gateway to AI Mastery"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = 'Segoe UI'
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Mission cards
        missions = [
            ("🎯", "Mission", "Empower UST students with\ncutting-edge AI skills"),
            ("🌟", "Vision", "Transform UST into Minnesota's\npremier AI talent hub"),
            ("🤝", "Community", "Where curiosity\nmeets capability"),
            ("🚀", "Promise", "Real skills, real projects,\nreal results")
        ]
        
        self.create_card_grid(slide, missions, start_y=Inches(3.2))
    
    def create_pillars_slide(self):
        """Create the four pillars slide with stunning layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Multi-color gradient
        self.add_gradient_background(slide, RGBColor(64, 93, 230), RGBColor(132, 94, 194))
        
        self.add_stunning_title(slide, "FOUR PILLARS OF AI EXCELLENCE", RGBColor(255, 255, 255))
        
        # Create 2x2 grid of pillars
        pillars = [
            ("🛠️", "HANDS-ON LEARNING", "Weekly coding workshops\nML model building sessions\nAI tool masterclasses"),
            ("🚀", "REAL PROJECTS", "Industry partnerships\nStartup collaborations\nSocial impact initiatives"),
            ("🏆", "COMPETITIONS", "Hackathons & contests\nResearch paper submissions\nInnovation challenges"),
            ("🌐", "NETWORKING", "Industry guest speakers\nAlumni mentorship\nCareer placement support")
        ]
        
        positions = [
            (Inches(1), Inches(3)),    # Top left
            (Inches(8.5), Inches(3)),  # Top right
            (Inches(1), Inches(5.5)),  # Bottom left
            (Inches(8.5), Inches(5.5)) # Bottom right
        ]
        
        for i, (icon, title, content) in enumerate(pillars):
            x, y = positions[i]
            
            # Create pillar card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, Inches(6), Inches(2.3)
            )
            
            # Gradient fill for each card
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            fill.transparency = 0.85
            
            # Add shadow
            shadow = card.shadow
            shadow.inherit = False
            shadow.style = 1
            shadow.blur_radius = Pt(15)
            shadow.distance = Pt(8)
            
            # Add icon
            icon_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(0.1), Inches(1), Inches(0.8))
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(48)
            icon_para.alignment = PP_ALIGN.CENTER
            
            # Add title
            title_box = slide.shapes.add_textbox(x + Inches(1.8), y + Inches(0.1), Inches(4), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Segoe UI'
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = self.NEXUS_BLUE
            
            # Add content
            content_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1), Inches(5), Inches(1.2))
            content_frame = content_box.text_frame
            content_frame.text = content
            content_para = content_frame.paragraphs[0]
            content_para.font.name = 'Segoe UI'
            content_para.font.size = Pt(12)
            content_para.font.color.rgb = self.NEXUS_DARK
    
    def create_card_grid(self, slide, items, start_y, cols=2):
        """Create a grid of beautiful cards"""
        rows = math.ceil(len(items) / cols)
        card_width = Inches(6.5)
        card_height = Inches(1.8)
        
        for i, (icon, title, content) in enumerate(items):
            row = i // cols
            col = i % cols
            
            x = Inches(1.5) + (col * Inches(7.5))
            y = start_y + (row * Inches(2.2))
            
            # Create card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            
            # Beautiful gradient fill
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            fill.transparency = 0.9
            
            # Add elegant shadow
            shadow = card.shadow
            shadow.inherit = False
            shadow.style = 1
            shadow.blur_radius = Pt(12)
            shadow.distance = Pt(6)
            shadow.angle = 135
            
            # Add icon
            icon_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.8))
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(36)
            icon_para.alignment = PP_ALIGN.CENTER
            
            # Add title
            title_box = slide.shapes.add_textbox(x + Inches(1.5), y + Inches(0.1), Inches(4.8), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Segoe UI'
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = self.NEXUS_BLUE
            
            # Add content
            content_box = slide.shapes.add_textbox(x + Inches(1.5), y + Inches(0.8), Inches(4.8), Inches(0.9))
            content_frame = content_box.text_frame
            content_frame.text = content
            content_para = content_frame.paragraphs[0]
            content_para.font.name = 'Segoe UI'
            content_para.font.size = Pt(11)
            content_para.font.color.rgb = self.NEXUS_DARK
    
    def add_stunning_title(self, slide, title_text, color):
        """Add a stunning title with effects"""
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        
        title_para = title_frame.paragraphs[0]
        title_para.font.name = 'Segoe UI Black'
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = color
        title_para.alignment = PP_ALIGN.CENTER
    
    def create_membership_slide(self):
        """Create a beautiful membership tiers slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Purple gradient background
        self.add_gradient_background(slide, RGBColor(155, 81, 224), RGBColor(255, 110, 199))
        
        self.add_stunning_title(slide, "READY TO SHAPE THE FUTURE?", RGBColor(255, 255, 255))
        
        # Membership tiers
        tiers = [
            ("🌟", "FULL MEMBERSHIP", "$20/semester", "All workshops & events\nCompetition eligibility\nMentorship matching\nIndustry networking", self.NEXUS_ACCENT),
            ("⚡", "CASUAL MEMBER", "FREE", "Monthly workshops\nSocial events\nResource access\nCommunity support", self.NEXUS_SUCCESS),
            ("🔥", "LEADERSHIP TRACK", "Application Required", "Lead project teams\nOrganize events\nBuild resume\nExclusive opportunities", RGBColor(255, 107, 53))
        ]
        
        card_width = Inches(4.5)
        card_height = Inches(4.5)
        start_x = Inches(1.25)
        y_pos = Inches(2.5)
        
        for i, (icon, title, price, benefits, accent_color) in enumerate(tiers):
            x_pos = start_x + (i * Inches(5))
            
            # Create tier card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y_pos, card_width, card_height
            )
            
            # Beautiful card styling
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            fill.transparency = 0.9
            
            # Premium shadow for premium look
            shadow = card.shadow
            shadow.inherit = False
            shadow.style = 1
            shadow.blur_radius = Pt(20)
            shadow.distance = Pt(10)
            shadow.angle = 135
            
            # Add icon
            icon_box = slide.shapes.add_textbox(x_pos + Inches(0.5), y_pos + Inches(0.2), Inches(3.5), Inches(1))
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(60)
            icon_para.alignment = PP_ALIGN.CENTER
            
            # Add title
            title_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(1.2), Inches(4.1), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Segoe UI'
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = accent_color
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add price
            price_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(1.8), Inches(4.1), Inches(0.6))
            price_frame = price_box.text_frame
            price_frame.text = price
            price_para = price_frame.paragraphs[0]
            price_para.font.name = 'Segoe UI'
            price_para.font.size = Pt(18)
            price_para.font.bold = True
            price_para.font.color.rgb = self.NEXUS_BLUE
            price_para.alignment = PP_ALIGN.CENTER
            
            # Add benefits
            benefits_box = slide.shapes.add_textbox(x_pos + Inches(0.3), y_pos + Inches(2.5), Inches(3.9), Inches(1.8))
            benefits_frame = benefits_box.text_frame
            benefits_frame.text = benefits
            benefits_para = benefits_frame.paragraphs[0]
            benefits_para.font.name = 'Segoe UI'
            benefits_para.font.size = Pt(11)
            benefits_para.font.color.rgb = self.NEXUS_DARK
            benefits_para.alignment = PP_ALIGN.LEFT
    
    def create_final_cta_slide(self):
        """Create a powerful call-to-action slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Electric gradient for energy
        self.add_gradient_background(slide, RGBColor(255, 154, 0), RGBColor(255, 206, 84))
        
        # Large, bold title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "YOUR AI JOURNEY\nSTARTS NOW"
        title_para = title_frame.paragraphs[0]
        title_para.font.name = 'Segoe UI Black'
        title_para.font.size = Pt(64)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Action steps
        actions_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(3))
        actions_frame = actions_box.text_frame
        actions_text = """🚀 IMMEDIATE ACTIONS:
1. Join Our Discord - Get connected today
2. Follow @NexusAI_UST - Stay updated
3. Attend Next Meeting - [Date/Time/Location]
4. Sign Up for Newsletter - Weekly opportunities

📞 Contact: nexusai@stthomas.edu | [Phone] | nexusai.stthomas.edu"""
        
        actions_frame.text = actions_text
        actions_para = actions_frame.paragraphs[0]
        actions_para.font.name = 'Segoe UI'
        actions_para.font.size = Pt(20)
        actions_para.font.color.rgb = RGBColor(255, 255, 255)
        actions_para.alignment = PP_ALIGN.CENTER
        
        # Add call-to-action button effect
        cta_button = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.5), Inches(7.2), Inches(5), Inches(1)
        )
        
        fill = cta_button.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Button shadow
        shadow = cta_button.shadow
        shadow.inherit = False
        shadow.style = 1
        shadow.blur_radius = Pt(15)
        shadow.distance = Pt(8)
        
        # Button text
        button_text = slide.shapes.add_textbox(Inches(5.5), Inches(7.4), Inches(5), Inches(0.6))
        button_frame = button_text.text_frame
        button_frame.text = "JOIN NEXUS AI TODAY!"
        button_para = button_frame.paragraphs[0]
        button_para.font.name = 'Segoe UI Black'
        button_para.font.size = Pt(24)
        button_para.font.bold = True
        button_para.font.color.rgb = self.NEXUS_BLUE
        button_para.alignment = PP_ALIGN.CENTER
    
    def create_all_slides(self):
        """Create the complete beautiful presentation"""
        print("🎨 Creating stunning title slide...")
        self.create_title_slide()
        
        print("🚀 Creating opportunity slide...")
        self.create_opportunity_slide()
        
        print("⚠️ Creating problem slide...")
        self.create_problem_slide()
        
        print("✨ Creating solution slide...")
        self.create_solution_slide()
        
        print("🏛️ Creating pillars slide...")
        self.create_pillars_slide()
        
        print("💎 Creating membership slide...")
        self.create_membership_slide()
        
        print("🎯 Creating call-to-action slide...")
        self.create_final_cta_slide()
        
        # Save the beautiful presentation
        output_file = "/Users/cloudaistudio/Documents/UST/NEXUS/Nexus_AI_BEAUTIFUL_Presentation.pptx"
        self.prs.save(output_file)
        print(f"✨ Beautiful presentation saved: {output_file}")
        
        return output_file

if __name__ == "__main__":
    print("🎨 Creating Beautiful Nexus AI Club Presentation...")
    print("🌟 This will look AMAZING!")
    
    creator = BeautifulNexusPresentation()
    output_path = creator.create_all_slides()
    
    print(f"\n🎉 SUCCESS! Beautiful presentation created!")
    print(f"📁 File: {output_path}")
    print(f"\n✨ Features included:")
    print(f"   • Stunning gradient backgrounds")
    print(f"   • Professional card layouts") 
    print(f"   • Beautiful typography")
    print(f"   • Drop shadows and effects")
    print(f"   • Modern color scheme")
    print(f"   • Perfect 16:9 aspect ratio")
    print(f"\n🚀 This presentation will WOW your audience!")

